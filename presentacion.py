from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_COLOR = RGBColor(15, 23, 42)
TEXT_WHITE = RGBColor(248, 250, 252)
ACCENT_CYAN = RGBColor(6, 182, 212)
ACCENT_GREEN = RGBColor(16, 185, 129)
TEXT_MUTED = RGBColor(148, 163, 184)

def aplicar_fondo(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

slides_data = [
    {
        "titulo": "SISTEMA ELECTORAL BLINDADO",
        "subtitulo": "Democracia Digital Inviolable | Cero Rastro\n\n• Arquitectura Zero-Trust   • Sellado SHA-256   • Doble Factor 2FA TOTP",
        "es_portada": True
    },
    {
        "titulo": "El Desafío del Voto Electrónico Convencional",
        "items": [
            "• El Mito del Anonimato: Bases de datos comunes vinculan al elector con su voto.",
            "• Vulnerabilidad de Suplantación: Accesos estáticos o correos sin verificación 2FA real.",
            "• Falta de Prueba Matemática: El usuario no puede verificar si su papeleta fue alterada.",
            "\n\"La legitimidad de una elección no se promete: se demuestra con datos inmutables.\""
        ]
    },
    {
        "titulo": "Arquitectura Desacoplada: Secreto por Diseño",
        "items": [
            "• Censo & Autenticación (BD 1): Administra identidad, contraseñas Bcrypt y secretos TOTP. Emite un token ciego efímero.",
            "• Urna Digital (BD 2): Recibe el voto anónimo y lo sella en la cadena SHA-256 sin conocer la identidad del votante.",
            "• Cero Relación Foránea: Imposibilidad técnica y matemática de reconstruir quién votó por quién."
        ]
    },
    {
        "titulo": "Flujo de Votación en 3 Pasos",
        "items": [
            "1. Autenticación Fuerte: Ingreso con Cédula + Código OTP dinámico (Authenticator).",
            "2. Cabina Limpia: Selección rápida e intuitiva del candidato en interfaz táctil.",
            "3. Sellado Criptográfico: Emisión inmediata del comprobante SHA-256 (64 caracteres)."
        ]
    },
    {
        "titulo": "El Verificador Criptográfico Independiente",
        "items": [
            "• Transparencia Activa: Cualquier elector o veedor puede consultar la inclusión de su papeleta.",
            "• Consulta Pública: Validación instantánea en la urna sin requerir inicio de sesión.",
            "• Inmutabilidad Garantizada: Si se intenta modificar un voto, se invalida toda la cadena pública."
        ]
    },
    {
        "titulo": "Control Institucional & Auditoría",
        "items": [
            "• Panel de Administración: Apertura con Clave de Custodia y control de personal.",
            "• Panel de Auditoría: Bitácora inmutable de eventos y visualización de escrutinio en vivo.",
            "• Acta Oficial de Cierre: Consolidado digital listo para descarga y firma de jurados."
        ]
    },
    {
        "titulo": "Infraestructura Cloud de Grado Empresarial",
        "items": [
            "• Frontend: React + Vite + Tailwind CSS desplegado en Vercel (Edge Global).",
            "• Backend: Node.js + TypeScript + Express en Render con RBAC estricto.",
            "• Persistencia: PostgreSQL desacoplado en Supabase con políticas RLS activas."
        ]
    },
    {
        "titulo": "¿Por qué Implementar este Sistema?",
        "items": [
            "• Cero Impugnaciones: Resultados verificables matemáticamente por todas las partes.",
            "• Ahorro Operativo: Eliminación total de gastos en tarjetones y papelería.",
            "• Escrutinio Instantáneo: Resultados definitivos disponibles en segundos tras el cierre."
        ]
    }
]

blank_layout = prs.slide_layouts[6]

for data in slides_data:
    slide = prs.slides.add_slide(blank_layout)
    aplicar_fondo(slide)
    
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.33), Inches(1.2))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data["titulo"]
    p.font.size = Pt(32 if not data.get("es_portada") else 40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(4.5))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    
    if data.get("es_portada"):
        cp = ctf.paragraphs[0]
        cp.text = data["subtitulo"]
        cp.font.size = Pt(20)
        cp.font.color.rgb = TEXT_WHITE
    else:
        for idx, item in enumerate(data.get("items", [])):
            cp = ctf.add_paragraph() if idx > 0 else ctf.paragraphs[0]
            cp.text = item
            cp.font.size = Pt(18)
            cp.space_after = Pt(14)
            cp.font.color.rgb = ACCENT_GREEN if item.startswith("•") or item[0].isdigit() else TEXT_MUTED

prs.save("Presentacion_Sistema_Electoral.pptx")
print("Presentación generada con éxito: Presentacion_Sistema_Electoral.pptx")